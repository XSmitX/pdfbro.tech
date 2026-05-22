"""PDF to PowerPoint converter using python-pptx and PyMuPDF.

Usage:
  python pdf_to_pptx.py <input.pdf> <output.pptx>

Exit codes:
  0 - success (prints SUCCESS:<output_path>)
  1 - error (message printed to stderr)
"""

import io
import os
import sys


def pdf_to_pptx(input_path: str, output_path: str) -> None:
    """Convert each PDF page into a full-slide image in a PPTX file."""
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation

        print("Converting PDF pages to images...", file=sys.stderr)
        doc = fitz.open(input_path)

        if doc.page_count == 0:
            raise RuntimeError("The input PDF has no pages")

        print("Creating PowerPoint presentation...", file=sys.stderr)
        prs = Presentation()

        # Keep 16:9 output to match the UI expectation.
        prs.slide_height = 9 * 914400
        prs.slide_width = 16 * 914400

        # 150 DPI gives decent quality without extreme file size.
        zoom = 150 / 72
        matrix = fitz.Matrix(zoom, zoom)

        total_pages = doc.page_count
        for i, page in enumerate(doc):
            print(f"Processing slide {i + 1}/{total_pages}...", file=sys.stderr)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            image_stream = io.BytesIO(pix.tobytes("png"))

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                image_stream,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        doc.close()
        prs.save(output_path)
        print(f"Successfully converted {total_pages} pages to PowerPoint.", file=sys.stderr)

    except ImportError as e:
        missing = str(e).split("'")[1] if "'" in str(e) else str(e)
        raise RuntimeError(
            f"Required library not installed: {missing}. "
            "Run: pip install python-pptx pymupdf"
        )
    except Exception as e:
        raise RuntimeError(f"Conversion failed: {e}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pdf_to_pptx.py <input.pdf> <output.pptx>", file=sys.stderr)
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    if not os.path.exists(input_file):
        print(f"Input file not found: {input_file}", file=sys.stderr)
        sys.exit(1)
    
    try:
        pdf_to_pptx(input_file, output_file)
        if os.path.exists(output_file):
            print(f"SUCCESS:{output_file}")
            sys.exit(0)
        else:
            print("Conversion produced no output file", file=sys.stderr)
            sys.exit(1)
    except Exception as e:
        print(str(e), file=sys.stderr)
        sys.exit(1)
